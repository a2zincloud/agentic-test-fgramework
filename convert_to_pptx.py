#!/usr/bin/env python3
"""
Convert Markdown Pitch Deck to PowerPoint Presentation
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

def create_presentation():
    """Create a professional PowerPoint presentation from the pitch deck"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (professional blue theme)
    TITLE_COLOR = RGBColor(26, 35, 126)  # Dark blue
    ACCENT_COLOR = RGBColor(74, 144, 226)  # Light blue
    TEXT_COLOR = RGBColor(33, 33, 33)  # Dark gray
    SUCCESS_COLOR = RGBColor(76, 175, 80)  # Green
    
    def add_title_slide(title, subtitle):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = TITLE_COLOR
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = ACCENT_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Tagline
        tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.8))
        tagline_frame = tagline_box.text_frame
        tagline_frame.text = "Reducing Testing Costs by 80% While Improving Quality"
        tagline_para = tagline_frame.paragraphs[0]
        tagline_para.font.size = Pt(24)
        tagline_para.font.italic = True
        tagline_para.font.color.rgb = TEXT_COLOR
        tagline_para.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_items, layout_type="bullet"):
        """Add content slide with bullets or table"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = TITLE_COLOR
        
        # Content area
        if layout_type == "bullet":
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for item in content_items:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = TEXT_COLOR
                p.space_before = Pt(12)
        
        elif layout_type == "two_column":
            # Left column
            left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4), Inches(5.5))
            left_frame = left_box.text_frame
            left_frame.word_wrap = True
            
            # Right column
            right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(5.5))
            right_frame = right_box.text_frame
            right_frame.word_wrap = True
            
            mid = len(content_items) // 2
            for item in content_items[:mid]:
                p = left_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_COLOR
                p.space_before = Pt(10)
            
            for item in content_items[mid:]:
                p = right_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_COLOR
                p.space_before = Pt(10)
    
    def add_table_slide(title, headers, rows):
        """Add slide with table"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = TITLE_COLOR
        
        # Table
        rows_count = len(rows) + 1  # +1 for header
        cols_count = len(headers)
        
        left = Inches(0.8)
        top = Inches(1.5)
        width = Inches(8.4)
        height = Inches(5)
        
        table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
        
        # Set column widths
        for i in range(cols_count):
            table.columns[i].width = Inches(8.4 / cols_count)
        
        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_COLOR
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        
        # Data rows
        for i, row in enumerate(rows):
            for j, cell_text in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(cell_text)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = TEXT_COLOR
    
    def add_big_number_slide(title, numbers):
        """Add slide with big numbers/metrics"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = TITLE_COLOR
        
        # Numbers in grid
        num_items = len(numbers)
        cols = 2
        rows = (num_items + 1) // 2
        
        for idx, (number, label) in enumerate(numbers):
            row = idx // cols
            col = idx % cols
            
            x = Inches(1 + col * 4.2)
            y = Inches(2 + row * 2)
            
            # Number
            num_box = slide.shapes.add_textbox(x, y, Inches(3.5), Inches(0.8))
            num_frame = num_box.text_frame
            num_frame.text = number
            num_para = num_frame.paragraphs[0]
            num_para.font.size = Pt(48)
            num_para.font.bold = True
            num_para.font.color.rgb = SUCCESS_COLOR
            num_para.alignment = PP_ALIGN.CENTER
            
            # Label
            label_box = slide.shapes.add_textbox(x, y + Inches(0.9), Inches(3.5), Inches(0.6))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(16)
            label_para.font.color.rgb = TEXT_COLOR
            label_para.alignment = PP_ALIGN.CENTER
    
    # Slide 1: Title
    add_title_slide(
        "🚀 Agentic Testing Framework",
        "Transforming Health Insurance Testing with AI"
    )
    
    # Slide 2: The Problem
    add_content_slide(
        "💥 The Testing Crisis in Health Insurance",
        [
            "Manual Testing Bottleneck: 60% of release cycle time ($2M+ annually)",
            "Complex Regulatory Requirements: HIPAA, ACA, State mandates ($500K+)",
            "Legacy System Integration: Mainframe, EDI, HL7 ($800K+)",
            "Production Defects: Customer impact, penalties ($1M+)",
            "Test Maintenance: 40% of QA time ($600K+)",
            "",
            "Reality: 20+ QA Engineers, 6-8 Week testing cycles",
            "30-40% of defects escape to production",
            "Total annual testing costs: $4.9M+"
        ]
    )
    
    # Slide 3: The Vision
    add_content_slide(
        "🎯 Imagine a World Where...",
        [
            "✨ Testing is Autonomous - AI agents write, execute, and maintain tests",
            "🔮 Quality is Predictive - Defects caught before they're written",
            "✅ Compliance is Built-In - Regulatory requirements validated real-time",
            "💰 Costs are Minimal - 80% reduction in testing expenses",
            "⚡ Speed is Exponential - Hours instead of weeks for full regression"
        ]
    )
    
    # Slide 4: The Solution
    add_content_slide(
        "🤖 Agentic Testing Framework",
        [
            "AI-Powered, Domain-Specific, Enterprise-Ready",
            "",
            "🧠 AI Orchestrator (GPT-4 / Claude / Custom Models)",
            "↓",
            "🎭 Specialized AI Agents:",
            "  • Policy Validation Agent",
            "  • Claims Processing Agent",
            "  • Member Journey Agent",
            "  • Integration Testing Agent",
            "  • Compliance & Security Agent",
            "↓",
            "🛠️ Testing Tools & Systems",
            "  • UI • API • Database • EDI • Compliance"
        ]
    )
    
    # Slide 5: How It Works
    add_content_slide(
        "⚙️ The Magic Behind the Scenes",
        [
            "1️⃣ Intelligent Test Generation",
            "   Policy Document → AI Analysis → 1000+ Test Cases",
            "   (Manual: 2 weeks vs AI: 2 hours)",
            "",
            "2️⃣ Autonomous Execution",
            "   • Parallel execution across 100+ agents",
            "   • Self-healing when UI/API changes",
            "   • Real-time adaptation to failures",
            "",
            "3️⃣ Smart Validation",
            "   • Claims adjudication logic verification",
            "   • Regulatory compliance checks",
            "   • Cross-system integration validation",
            "",
            "4️⃣ Continuous Learning",
            "   • Learns from production issues",
            "   • Improves test coverage over time"
        ]
    )
    
    # Slide 6: Domain Expertise
    add_content_slide(
        "🏥 Built for Health Insurance",
        [
            "Policy Management:",
            "  • Premium calculations • Eligibility rules",
            "  • Coverage limits • Renewal workflows",
            "",
            "Claims Processing:",
            "  • Adjudication logic • COB scenarios",
            "  • Pre-authorization • Payment accuracy",
            "",
            "Compliance:",
            "  • HIPAA validation • ACA requirements",
            "  • State mandates • Audit trails",
            "",
            "Integration:",
            "  • EDI (837, 835, 834, 270/271)",
            "  • HL7/FHIR • Provider networks • PBM systems"
        ],
        layout_type="two_column"
    )
    
    # Slide 7: Competitive Advantage
    add_table_slide(
        "🏆 Why We Win",
        ["Feature", "Manual", "Traditional", "Our Framework"],
        [
            ["Speed", "6-8 weeks", "2-3 weeks", "2-3 days ✨"],
            ["Coverage", "40-50%", "60-70%", "90%+ ✨"],
            ["Maintenance", "High", "High", "Low ✨"],
            ["Domain Knowledge", "Manual", "Manual", "Built-in AI ✨"],
            ["Cost (3 years)", "$6M", "$3M", "$1.5M ✨"],
            ["Adaptability", "Medium", "Low", "High ✨"]
        ]
    )
    
    # Slide 8: ROI
    add_big_number_slide(
        "💰 The Numbers That Matter",
        [
            ("145-217%", "Year 1 ROI"),
            ("$2.5M-$3.8M", "Annual Savings"),
            ("6-8 Months", "Payback Period"),
            ("80%", "Cost Reduction")
        ]
    )
    
    # Slide 9: Business Impact
    add_content_slide(
        "📈 Transformational Outcomes",
        [
            "Efficiency Gains:",
            "  ⚡ 80% reduction in test creation time",
            "  🚀 70% reduction in test execution time",
            "  🔧 90% reduction in test maintenance",
            "  📊 3x increase in test coverage",
            "",
            "Quality Improvements:",
            "  🎯 95%+ defect detection in pre-production",
            "  📉 50% reduction in production defects",
            "  ✅ 98%+ test reliability",
            "  🔍 100% compliance validation",
            "",
            "Business Benefits:",
            "  💵 $2.5M+ annual cost savings",
            "  ⏱️ 30% faster time to market"
        ]
    )
    
    # Slide 10: Enterprise Ready
    add_content_slide(
        "🏢 Built for Large Insurance Companies",
        [
            "Security & Compliance:",
            "  ✅ HIPAA compliant with PHI/PII protection",
            "  ✅ SOC 2 Type II certified",
            "  ✅ HITRUST CSF framework",
            "  ✅ Zero-trust architecture",
            "",
            "Scalability:",
            "  ✅ 10,000+ concurrent tests",
            "  ✅ Multi-region deployment",
            "  ✅ 99.9% uptime SLA",
            "",
            "Integration:",
            "  ✅ Legacy systems (Mainframe, AS/400)",
            "  ✅ Modern APIs (REST, GraphQL, gRPC)",
            "  ✅ Enterprise tools (JIRA, ServiceNow)"
        ]
    )
    
    # Slide 11: Implementation Roadmap
    add_content_slide(
        "🗺️ 12-Month Journey to Success",
        [
            "Phase 1: Foundation (Months 1-3)",
            "  • Core framework deployment",
            "  • 100+ automated tests",
            "  • Milestone: 50% reduction in manual testing",
            "",
            "Phase 2: Scale (Months 4-6)",
            "  • All 5 specialized agents",
            "  • 1,000+ automated tests",
            "  • Milestone: 70% test automation coverage",
            "",
            "Phase 3: Advanced (Months 7-9)",
            "  • AI-powered test generation",
            "  • 5,000+ automated tests",
            "  • Milestone: <1 hour regression",
            "",
            "Phase 4: Optimization (Months 10-12)",
            "  • Multi-region deployment",
            "  • Milestone: 80% automation, ROI positive"
        ]
    )
    
    # Slide 12: Success Metrics
    add_big_number_slide(
        "🌟 What Success Looks Like",
        [
            ("90%+", "Test Coverage"),
            ("3 days", "Full Regression"),
            ("$2.1M", "Cost Savings"),
            ("40%", "Faster Releases")
        ]
    )
    
    # Slide 13: Investment Ask
    add_content_slide(
        "💼 The Ask",
        [
            "Pilot Program Investment: $250,000 (3 months)",
            "  • Framework setup: $100K",
            "  • 2 specialized agents: $80K",
            "  • Integration: $40K",
            "  • Training & support: $30K",
            "",
            "Expected Pilot Outcomes:",
            "  ✅ 100+ automated tests",
            "  ✅ 50% reduction in manual testing",
            "  ✅ Proof of ROI",
            "  ✅ Executive buy-in for full deployment",
            "",
            "Full Program: $1.15M - $1.75M (12 months)",
            "Expected Return: $2.5M - $3.8M annually",
            "ROI: 145% - 217% in Year 1"
        ]
    )
    
    # Slide 14: Why Now
    add_content_slide(
        "⏰ The Time is Now",
        [
            "Market Trends:",
            "  📈 85% of enterprises investing in AI",
            "  🚀 60% cite testing as bottleneck",
            "  💰 Need to do more with less",
            "",
            "Competitive Pressure:",
            "  • Competitors exploring AI testing",
            "  • First movers gain significant advantage",
            "  • Technology gap widening rapidly",
            "",
            "Your Opportunity:",
            "  🎯 Be a pioneer in AI-powered testing",
            "  🏆 Gain competitive advantage",
            "  💡 Transform from cost center to innovation driver",
            "  🚀 Lead the industry in testing excellence"
        ]
    )
    
    # Slide 15: Call to Action
    add_content_slide(
        "🎯 Let's Transform Testing Together",
        [
            "Next Steps:",
            "",
            "1. Pilot Program (3 months) - $250K",
            "   • Select one product line",
            "   • Deploy core framework",
            "   • Demonstrate quick wins",
            "",
            "2. Business Case Approval (Month 4)",
            "   • Present results to executives",
            "   • Secure full budget",
            "",
            "3. Full Deployment (9 months) - $900K-$1.5M",
            "   • Phased rollout across products",
            "   • Scale to enterprise",
            "   • Achieve target ROI",
            "",
            "Decision Timeline: 4 weeks to pilot kickoff"
        ]
    )
    
    # Slide 16: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🏗️ System Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Architecture layers
    layers = [
        ("User Interface & API Gateway", 1.5, RGBColor(100, 181, 246)),
        ("Orchestration Layer (LLM-based)", 2.4, RGBColor(74, 144, 226)),
        ("Agent Layer (5 Specialized Agents)", 3.3, RGBColor(66, 165, 245)),
        ("Tool Layer (UI, API, DB, EDI, Reports)", 4.2, RGBColor(41, 128, 185)),
        ("Data Layer (Vector DB, PostgreSQL, Redis)", 5.1, RGBColor(52, 152, 219))
    ]
    
    for layer_name, y_pos, color in layers:
        # Layer box
        box = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1.5), Inches(y_pos),
            Inches(7), Inches(0.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(2)
        
        # Layer text
        text_frame = box.text_frame
        text_frame.text = layer_name
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add arrows between layers
    for i in range(len(layers) - 1):
        arrow = slide.shapes.add_shape(
            5,  # Down arrow
            Inches(4.8), Inches(layers[i][1] + 0.75),
            Inches(0.4), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
        arrow.line.color.rgb = RGBColor(100, 100, 100)
    
    # Slide 17: Agent Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🤖 Specialized Agent Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Central orchestrator
    center_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(3.5), Inches(1.8),
        Inches(3), Inches(0.8)
    )
    center_box.fill.solid()
    center_box.fill.fore_color.rgb = RGBColor(255, 152, 0)
    center_box.line.color.rgb = RGBColor(230, 126, 34)
    center_box.line.width = Pt(3)
    
    text_frame = center_box.text_frame
    text_frame.text = "AI Orchestrator\n(GPT-4/Claude)"
    text_frame.paragraphs[0].font.size = Pt(16)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Specialized agents in circle
    agents = [
        ("Policy\nAgent", 1.5, 3.5, RGBColor(76, 175, 80)),
        ("Claims\nAgent", 3.5, 3.0, RGBColor(33, 150, 243)),
        ("Member\nAgent", 5.5, 3.5, RGBColor(156, 39, 176)),
        ("Integration\nAgent", 5.5, 5.0, RGBColor(255, 87, 34)),
        ("Security\nAgent", 3.5, 5.5, RGBColor(244, 67, 54)),
        ("Analytics\nAgent", 1.5, 5.0, RGBColor(0, 150, 136))
    ]
    
    for agent_name, x_pos, y_pos, color in agents:
        agent_box = slide.shapes.add_shape(
            9,  # Rounded rectangle
            Inches(x_pos), Inches(y_pos),
            Inches(1.5), Inches(0.8)
        )
        agent_box.fill.solid()
        agent_box.fill.fore_color.rgb = color
        agent_box.line.color.rgb = RGBColor(255, 255, 255)
        agent_box.line.width = Pt(2)
        
        text_frame = agent_box.text_frame
        text_frame.text = agent_name
        text_frame.paragraphs[0].font.size = Pt(13)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Slide 18: Data Flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📊 Test Execution Flow"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Flow steps
    flow_steps = [
        ("1. User\nRequest", RGBColor(66, 165, 245)),
        ("2. Orchestrator\nAnalysis", RGBColor(74, 144, 226)),
        ("3. Agent\nSelection", RGBColor(41, 128, 185)),
        ("4. Test\nGeneration", RGBColor(52, 152, 219)),
        ("5. Execution", RGBColor(30, 136, 229)),
        ("6. Validation", RGBColor(25, 118, 210)),
        ("7. Results", RGBColor(21, 101, 192))
    ]
    
    x_start = 1.2
    for i, (step_name, color) in enumerate(flow_steps):
        step_box = slide.shapes.add_shape(
            9,  # Rounded rectangle
            Inches(x_start + i * 1.1), Inches(2.8),
            Inches(1), Inches(1.2)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = color
        step_box.line.color.rgb = RGBColor(255, 255, 255)
        step_box.line.width = Pt(2)
        
        text_frame = step_box.text_frame
        text_frame.text = step_name
        text_frame.paragraphs[0].font.size = Pt(11)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Add arrow
        if i < len(flow_steps) - 1:
            arrow = slide.shapes.add_shape(
                13,  # Right arrow
                Inches(x_start + i * 1.1 + 1.05), Inches(3.3),
                Inches(0.35), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
            arrow.line.color.rgb = RGBColor(100, 100, 100)
    
    # Slide 19: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_COLOR
    
    # Thank you text
    thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "🙏 Thank You"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(60)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = RGBColor(255, 255, 255)
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Let's Revolutionize Health Insurance Testing Together"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('Agentic_Testing_Framework_Pitch_Deck.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📄 File: Agentic_Testing_Framework_Pitch_Deck.pptx")
    print("📊 Total slides: 19 (including 3 architecture diagrams)")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx library not found")
        print("📦 Install it using: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")

